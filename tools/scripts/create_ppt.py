#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI Coding Showcase PPT Generator
生成包含全量设计理念的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 颜色主题
TITLE_BLUE = RGBColor(0, 51, 102)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(102, 102, 102)
ACCENT_RED = RGBColor(153, 51, 51)
ACCENT_GREEN = RGBColor(0, 102, 51)
ACCENT_ORANGE = RGBColor(204, 102, 0)
ACCENT_PURPLE = RGBColor(102, 51, 153)

def add_title_slide(prs, title, subtitle=""):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(102, 102, 102)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items, two_column=False):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # 内容区域
    if two_column and len(content_items) > 1:
        # 左栏
        left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(4.6), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True

        mid = len(content_items[0]) // 2 if len(content_items[0]) > 2 else len(content_items[0])

        for i, item in enumerate(content_items[0][:mid]):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            add_formatted_text(p, item)

        # 右栏
        right_box = slide.shapes.add_textbox(Inches(5.1), Inches(1.1), Inches(4.6), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items[0][mid:]):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            add_formatted_text(p, item)
    else:
        # 单栏
        content_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.1), Inches(9.4), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            add_formatted_text(p, item)

    return slide

def add_formatted_text(p, item):
    """添加格式化文本"""
    if isinstance(item, tuple):
        text, level, is_bold = item
        p.text = text
        p.level = level
        p.font.size = Pt(18 - level * 2)
        p.font.bold = is_bold
        if level == 0:
            p.font.color.rgb = RGBColor(0, 51, 102)
        else:
            p.font.color.rgb = RGBColor(51, 51, 51)
    else:
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加对比幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # 左侧标题
    left_title_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(153, 51, 51)  # 红色

    # 左侧内容
    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        add_formatted_text(p, item)

    # 右侧标题
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 51)  # 绿色

    # 右侧内容
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        add_formatted_text(p, item)

    return slide

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ===== Slide 1: 封面 =====
add_title_slide(prs, "AI Coding 探索", "基于 Claude Code 的 Agent 开发实践")

# ===== Slide 2: 研发痛点 =====
add_content_slide(prs, "研发痛点问题", [
    ("痛点 1: 上下文理解成本高", 0, True),
    ("每次对话需重新解释项目结构", 1, False),
    ("AI 缺乏项目知识积累机制", 1, False),
    ("痛点 2: 开发流程碎片化", 0, True),
    ("构建/测试/部署需切换多个环境", 1, False),
    ("错误修复缺乏系统性方法", 1, False),
    ("跨平台开发 (Windows 编辑 + Linux 运行) 复杂", 1, False),
    ("痛点 3: Token 成本指数级增长", 0, True),
    ("连续对话导致历史上下文堆积", 1, False),
    ("无效信息被重复加载", 1, False),
])

# ===== Slide 3: Agent 架构理念 =====
add_content_slide(prs, "Google Agent 白皮书核心理念", [
    ("Agent = Model + Tools + Orchestration + Knowledge", 0, True),
    ("", 0, False),
    ("Model (大模型)", 0, True),
    ("推理能力 | 语言理解 | 决策制定", 1, False),
    ("Tools (工具集)", 0, True),
    ("文件操作 | 代码执行 | 网络请求", 1, False),
    ("Orchestration (编排层)", 0, True),
    ("规划 | 推理 | 决策", 1, False),
    ("Knowledge (知识库)", 0, True),
    ("项目记忆 | 最佳实践 | 错误修复", 1, False),
    ("", 0, False),
    ("工作循环: 感知 -> 推理 -> 行动 -> 观察 -> 循环", 0, True),
])

# ===== Slide 4: Skills 体系架构 =====
add_content_slide(prs, "项目解决方案 - Skills 体系", [
    ("CLAUDE.md (项目规范) -> Orchestration", 0, True),
    ("结构说明、开发指南、模块介绍", 1, False),
    ("Skills/ (技能定义) -> Tools", 0, True),
    ("dev, fix-*, code-*, openresty-lua-plugins", 1, False),
    ("Memory/ (知识积累) -> Knowledge", 0, True),
    ("已解决问题、架构决策、最佳实践", 1, False),
    ("", 0, False),
    ("分类:", 0, True),
    ("核心开发: dev, openresty-lua-plugins, web-admin-frontend", 1, False),
    ("错误修复: fix-compile, fix-test, fix-runtime, fix-loop", 1, False),
    ("代码质量: code-review, code-reactor, feedback", 1, False),
    ("DevOps: jenkins-pipeline, monitor-observability", 1, False),
])

# ===== Slide 5: 跨平台开发模式 =====
add_content_slide(prs, "设计理念 1: 跨平台开发模式", [
    ("Windows 本地开发 + Linux 远程运行", 0, True),
    ("", 0, False),
    ("Windows 本地环境:", 0, True),
    ("  - 代码编辑 (IDE/编辑器)", 1, False),
    ("  - Git 版本管理", 1, False),
    ("  - 通过 /dev CLI 控制远程", 1, False),
    ("", 0, False),
    ("Linux 远程环境:", 0, True),
    ("  - 所有 Shell 脚本实际执行", 1, False),
    ("  - OpenResty 构建/运行", 1, False),
    ("  - 测试/服务管理", 1, False),
    ("", 0, False),
    ("SSH (paramiko) 实现透明跨平台", 0, True),
])

# ===== Slide 6: Hooks 安全机制 =====
add_content_slide(prs, "设计理念 2: Hooks 安全机制", [
    ("三层 Hook 架构: Pre -> 执行 -> Post -> Audit", 0, True),
    ("", 0, False),
    ("PreToolUse.sh - 执行前检查:", 0, True),
    ("  - 危险命令拦截 (rm -rf /, dd, mkfs)", 1, False),
    ("  - Git 危险操作警告 (force push, reset --hard)", 1, False),
    ("  - 配置文件删除保护", 1, False),
    ("", 0, False),
    ("PostToolUse.sh - 执行后处理:", 0, True),
    ("  - 错误类型自动检测 (compile/test/runtime)", 1, False),
    ("  - 智能 Skill 建议 (/fix-compile, /fix-test)", 1, False),
    ("  - 失败状态持久化", 1, False),
    ("", 0, False),
    ("AuditLog.sh - 审计日志:", 0, True),
    ("  - JSONL 格式记录所有操作", 1, False),
    ("  - 支持回溯和分析", 1, False),
])

# ===== Slide 7: 自动修复循环 =====
add_content_slide(prs, "设计理念 3: 自动修复循环", [
    ("/fix-loop 自动迭代修复", 0, True),
    ("", 0, False),
    ("执行流程:", 0, True),
    ("  执行命令 -> 分析结果 -> AI修复 -> 重新执行 -> 循环", 1, False),
    ("", 0, False),
    ("错误检测规则 (tools/fixers/*.yaml):", 0, True),
    ("  - compile/*.yaml: 编译错误规则", 1, False),
    ("  - runtime/*.yaml: 运行时错误规则", 1, False),
    ("  - test/*.yaml: 测试失败规则", 1, False),
    ("", 0, False),
    ("终止条件:", 0, True),
    ("  - 执行成功", 1, False),
    ("  - 达到最大循环次数 (默认 5 次)", 1, False),
    ("  - 无法识别的错误类型", 1, False),
])

# ===== Slide 8: 反馈优化机制 =====
add_content_slide(prs, "设计理念 4: 反馈优化机制", [
    ("/feedback 自动学习并更新知识库", 0, True),
    ("", 0, False),
    ("触发条件:", 0, True),
    ("  - 用户说 '记住这个'", 1, False),
    ("  - 修复重复出现的错误", 1, False),
    ("  - 发现新的最佳实践", 1, False),
    ("", 0, False),
    ("反馈目标:", 0, True),
    ("  - MEMORY.md: 会话记忆、临时发现", 1, False),
    ("  - SKILL 文件: 技能知识、修复规则", 1, False),
    ("  - CLAUDE.md: 项目规范、工作流程", 1, False),
    ("  - patterns.yaml: 可复用模式库", 1, False),
    ("", 0, False),
    ("形成持续优化的闭环", 0, True),
])

# ===== Slide 9: Token 优化策略 =====
add_content_slide(prs, "设计理念 5: Token 优化策略", [
    ("核心思想: 细粒度拆分，按需加载", 0, True),
    ("", 0, False),
    (".claudeignore - 忽略构建产物、依赖库、密钥", 0, True),
    ("  减少 40% 无效加载", 1, False),
    ("", 0, False),
    ("CLAUDE.md - 结构化项目介绍", 0, True),
    ("  避免重复理解项目上下文", 1, False),
    ("", 0, False),
    ("/clear - 模块切换时清理历史", 0, True),
    ("  阻断 Token 指数增长", 1, False),
    ("", 0, False),
    ("Skill Reference 按需加载:", 0, True),
    ("  SKILL.md (技能定义) 始终加载", 1, False),
    ("  references/*.md (详细文档) 按需加载", 1, False),
    ("  实现零 Token 消耗执行", 1, False),
])

# ===== Slide 10: 端到端流程对比 =====
add_comparison_slide(prs, "端到端开发流程对比",
    "传统方式",
    [
        ("[查阅文档] 2-3 小时", 0, False),
        ("[编写代码] 1-2 小时", 0, False),
        ("[本地编译] 30 分钟", 0, False),
        ("[发现问题] 不确定", 0, False),
        ("", 0, False),
        ("总计: 4-6 小时", 0, True),
    ],
    "Agent 辅助方式",
    [
        ("/openresty-lua-plugins 5 分钟", 0, False),
        ("/dev all (sync->build->test) 10 分钟", 0, False),
        ("/fix-* (自动修复，如需要)", 0, False),
        ("", 0, False),
        ("总计: 15-30 分钟", 0, True),
        ("", 0, False),
        ("效率提升: 8-16 倍", 0, True),
    ]
)

# ===== Slide 11: /dev 命令体系 =====
add_content_slide(prs, "/dev 统一命令入口", [
    ("零 Token 消耗执行工程任务", 0, True),
    ("", 0, False),
    ("/dev all [module]      # 完整流水线 (sync -> build -> start -> test)", 1, False),
    ("/dev build [module]    # 编译指定模块或全部", 1, False),
    ("/dev test [module]     # 测试指定模块或全部", 1, False),
    ("/dev test --dt [file]  # 运行 Test::Nginx 测试用例", 1, False),
    ("/dev start             # 启动服务（远程）", 1, False),
    ("/dev stop              # 停止服务（远程）", 1, False),
    ("/dev sync              # 同步代码到远程服务器", 1, False),
    ("/dev status            # 查看项目状态", 1, False),
    ("/dev analyze <type>    # 分析错误输出", 1, False),
    ("", 0, False),
    ("失败时自动建议修复 Skill", 0, True),
])

# ===== Slide 12: Skill 示例 - openresty-lua-plugins =====
add_content_slide(prs, "Skill 示例: openresty-lua-plugins", [
    ("一键生成完整的 OpenResty Lua 插件", 0, True),
    ("", 0, False),
    ("设计原则:", 0, True),
    ("  1. 优先使用 Nginx 原生配置 (proxy_pass, upstream, limit_req)", 1, False),
    ("  2. 禁止高风险行为 (阻塞 I/O, 全局变量, 无限循环)", 1, False),
    ("  3. 使用非阻塞 API (ngx.socket, ngx.shared.DICT)", 1, False),
    ("", 0, False),
    ("生成产物:", 0, True),
    ("  - plugins/{name}/{name}.lua    # 插件主文件", 1, False),
    ("  - plugins/{name}/README.md     # 插件文档", 1, False),
    ("  - plugins/{name}/nginx.conf.example  # 示例配置", 1, False),
    ("  - test/dt/{name}/basic.t       # DT 测试用例", 1, False),
])

# ===== Slide 13: 架构图 =====
add_content_slide(prs, "整体架构图", [
    ("┌─────────────────────────────────────────────────────────────┐", 0, False),
    ("│                    Claude Code Agent                        │", 0, True),
    ("├─────────────────────────────────────────────────────────────┤", 0, False),
    ("│  CLAUDE.md          Memory/            Skills/              │", 0, False),
    ("│  (项目规范)          (知识积累)         (技能定义)            │", 0, False),
    ("├─────────────────────────────────────────────────────────────┤", 0, False),
    ("│                    Hooks 安全层                              │", 0, True),
    ("│  PreToolUse -> [执行] -> PostToolUse -> AuditLog             │", 0, False),
    ("├─────────────────────────────────────────────────────────────┤", 0, False),
    ("│  /dev CLI          /fix-*            /feedback              │", 0, False),
    ("│  (统一入口)         (自动修复)         (反馈优化)             │", 0, False),
    ("├─────────────────────────────────────────────────────────────┤", 0, False),
    ("│               SSH (paramiko) 跨平台桥接                       │", 0, True),
    ("├─────────────────────────────────────────────────────────────┤", 0, False),
    ("│  Windows 本地                          Linux 远程            │", 0, False),
    ("│  - 代码编辑                            - Shell 执行          │", 0, False),
    ("│  - Git 管理                            - OpenResty 运行      │", 0, False),
    ("└─────────────────────────────────────────────────────────────┘", 0, False),
])

# ===== Slide 14: 总结 =====
add_content_slide(prs, "总结: 五大设计理念", [
    ("1. 跨平台开发模式", 0, True),
    ("   Windows 开发 + Linux 远程运行，SSH 透明桥接", 1, False),
    ("", 0, False),
    ("2. Hooks 安全机制", 0, True),
    ("   Pre/Post/Audit 三层防护，智能错误检测和建议", 1, False),
    ("", 0, False),
    ("3. 自动修复循环", 0, True),
    ("   /fix-loop 迭代修复，YAML 规则匹配错误类型", 1, False),
    ("", 0, False),
    ("4. 反馈优化机制", 0, True),
    ("   /feedback 持续学习，更新知识库形成闭环", 1, False),
    ("", 0, False),
    ("5. Token 优化策略", 0, True),
    ("   渐进式加载、Reference 按需读取、零消耗执行", 1, False),
])

# ===== Slide 15: 效果 =====
add_content_slide(prs, "实施效果", [
    ("开发效率提升", 0, True),
    ("  传统 4-6 小时 -> Agent 辅助 15-30 分钟", 1, False),
    ("  效率提升: 8-16 倍", 1, False),
    ("", 0, False),
    ("Token 成本降低", 0, True),
    ("  .claudeignore 减少 40% 无效加载", 1, False),
    ("  Reference 按需加载避免冗余", 1, False),
    ("", 0, False),
    ("开发体验改善", 0, True),
    ("  一键端到端: /dev all", 1, False),
    ("  自动修复: /fix-loop", 1, False),
    ("  知识积累: /feedback", 1, False),
    ("", 0, False),
    ("范式跃迁: 人驱动工具 -> Agent 驱动开发", 0, True),
])

# ===== Slide 16: Q&A =====
slide = add_title_slide(prs, "Thank You", "Q & A")

# 添加演示命令
demo_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1))
tf = demo_box.text_frame
p = tf.paragraphs[0]
p.text = "演示: /openresty-lua-plugins -> /dev all -> /feedback"
p.font.size = Pt(18)
p.font.color.rgb = LIGHT_GRAY
p.alignment = PP_ALIGN.CENTER

# 保存
output_path = "D:/m30020610/WorkSpace/Skills/docs/ai-coding-showcase-optimized.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
print(f"Total slides: {len(prs.slides)}")