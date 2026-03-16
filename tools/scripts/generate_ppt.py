#!/usr/bin/env python3
"""
AI Coding Showcase PPT Generator
生成 AI Coding 探索 Showcase 演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色定义
DARK_BLUE = RGBColor(26, 26, 46)
CYAN = RGBColor(0, 212, 255)
YELLOW = RGBColor(255, 204, 0)
RED = RGBColor(255, 59, 48)
GREEN = RGBColor(52, 199, 89)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(136, 136, 136)


def add_title_slide(title, subtitle):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = YELLOW
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(title, content_list, subtitle=""):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY

    # 内容
    start_y = 1.6 if subtitle else 1.3
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(start_y), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

    return slide


def add_pain_point_slide():
    """添加痛点问题幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "研发痛点问题"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 痛点列表
    pain_points = [
        ("痛点 1: 上下文理解成本高", [
            "• 每次对话需要重新解释项目结构",
            "• AI 缺乏项目知识积累机制",
            "• 重复性沟通消耗大量时间"
        ]),
        ("痛点 2: 开发流程碎片化", [
            "• 构建/测试/部署 需要切换多个环境",
            "• 错误修复缺乏系统性方法",
            "• 跨平台开发（Windows编辑 + Linux运行）复杂"
        ]),
        ("痛点 3: Token 成本指数级增长", [
            "• 连续对话导致历史上下文堆积",
            "• 无效信息被重复加载",
            "• 长期对话成本难以控制"
        ])
    ]

    y_pos = 1.4
    for title_text, items in pain_points:
        # 痛点标题框
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(1.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(80, 30, 30)
        shape.line.color.rgb = RED

        # 痛点标题
        title_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(11.9), Inches(0.4))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RED

        # 痛点内容
        content_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.5), Inches(11.9), Inches(1.0))
        tf = content_shape.text_frame
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE

        y_pos += 1.9

    return slide


def add_architecture_slide():
    """添加架构幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Google Agent 白皮书核心理念"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 公式
    formula_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
    tf = formula_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Agent = Model + Tools + Orchestration"
    p.font.size = Pt(24)
    p.font.color.rgb = YELLOW
    p.alignment = PP_ALIGN.CENTER

    # 架构图 - Orchestration
    orch_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.8), Inches(4.333), Inches(1.0))
    orch_shape.fill.solid()
    orch_shape.fill.fore_color.rgb = RGBColor(50, 50, 80)
    orch_shape.line.color.rgb = YELLOW
    tf = orch_shape.text_frame
    tf.paragraphs[0].text = "Orchestration (编排层)"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = YELLOW
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "规划、推理、决策"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 连接线文字
    line_text = slide.shapes.add_textbox(Inches(6.4), Inches(2.85), Inches(0.5), Inches(0.3))
    tf = line_text.text_frame
    p = tf.paragraphs[0]
    p.text = "│"
    p.font.size = Pt(20)
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # 三个核心组件
    components = [
        ("Model (大模型)", ["推理能力", "语言理解", "决策制定"], Inches(0.8)),
        ("Tools (工具集)", ["文件操作", "代码执行", "网络请求"], Inches(4.9)),
        ("Knowledge (知识库)", ["项目记忆", "最佳实践", "错误修复"], Inches(9.0))
    ]

    for title, items, x_pos in components:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.3), Inches(3.5), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(40, 40, 70)
        shape.line.color.rgb = CYAN

        tf = shape.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CYAN
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT

    # 工作循环
    cycle_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.0))
    cycle_shape.fill.solid()
    cycle_shape.fill.fore_color.rgb = RGBColor(30, 60, 80)
    cycle_shape.line.color.rgb = CYAN

    tf = cycle_shape.text_frame
    tf.paragraphs[0].text = "Agent 工作循环:  感知 → 推理 → 行动 → 观察 → 循环"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_solution_slide():
    """添加解决方案幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "项目解决方案 - Skills 体系"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 架构映射
    mapping_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Inches(2.8))
    mapping_box.fill.solid()
    mapping_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
    mapping_box.line.color.rgb = CYAN

    # 映射内容
    mappings = [
        ("CLAUDE.md (项目规范)", "Orchestration", "结构说明、开发指南、模块介绍"),
        ("Skills/ (技能定义)", "Tools", "dev, fix-*, code-*, openresty-lua-plugins"),
        ("Memory/ (知识积累)", "Knowledge", "已解决问题、架构决策、最佳实践")
    ]

    y_pos = 1.4
    for project_part, agent_part, desc in mappings:
        # 项目部分
        part_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(3.5), Inches(0.7))
        tf = part_box.text_frame
        tf.paragraphs[0].text = project_part
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CYAN

        # 箭头
        arrow_box = slide.shapes.add_textbox(Inches(4.3), Inches(y_pos), Inches(0.5), Inches(0.7))
        tf = arrow_box.text_frame
        tf.paragraphs[0].text = "→"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = YELLOW

        # Agent部分
        agent_box = slide.shapes.add_textbox(Inches(4.9), Inches(y_pos), Inches(2.5), Inches(0.7))
        tf = agent_box.text_frame
        tf.paragraphs[0].text = agent_part
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = YELLOW

        # 描述
        desc_box = slide.shapes.add_textbox(Inches(7.5), Inches(y_pos), Inches(5.0), Inches(0.7))
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE

        y_pos += 0.8

    # Skills分类表
    table_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Skills 分类:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = YELLOW

    # 表格
    table_data = [
        ("核心开发", "dev, openresty-lua-plugins, web-admin-frontend", "开发入口、插件生成"),
        ("错误修复", "fix-compile, fix-test, fix-runtime, fix-loop", "自动分析修复"),
        ("代码质量", "code-review, code-reactor, feedback", "质量保障、知识沉淀"),
        ("DevOps", "jenkins-pipeline, monitor-observability", "CI/CD、监控运维")
    ]

    y_pos = 4.7
    for cat, skills, desc in table_data:
        row_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.5))
        row_box.fill.solid()
        row_box.fill.fore_color.rgb = RGBColor(50, 50, 80)
        row_box.line.color.rgb = RGBColor(80, 80, 100)

        # 类别
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.1), Inches(2.0), Inches(0.3))
        tf = cat_box.text_frame
        tf.paragraphs[0].text = cat
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CYAN

        # Skills
        skills_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos + 0.1), Inches(5.5), Inches(0.3))
        tf = skills_box.text_frame
        tf.paragraphs[0].text = skills
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE

        # 作用
        desc_box = slide.shapes.add_textbox(Inches(8.3), Inches(y_pos + 0.1), Inches(4.3), Inches(0.3))
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = GRAY

        y_pos += 0.55

    return slide


def add_comparison_slide():
    """添加对比幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "端到端开发流程实例"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 场景说明
    scene_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.5))
    tf = scene_box.text_frame
    p = tf.paragraphs[0]
    p.text = "场景：开发一个 OpenResty Lua 插件"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY

    # 传统方式
    old_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.0), Inches(3.5))
    old_box.fill.solid()
    old_box.fill.fore_color.rgb = RGBColor(80, 30, 30)
    old_box.line.color.rgb = RED

    old_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = old_title.text_frame
    tf.paragraphs[0].text = "传统方式"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED

    old_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(5.6), Inches(2.5))
    tf = old_content.text_frame
    tf.word_wrap = True
    items = ["[查阅文档] 2-3小时", "[编写代码] 1-2小时", "[本地编译] 30分钟", "[发现问题] 不确定时间", "", "总计: 4-6 小时"]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

    # Agent辅助方式
    new_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.6), Inches(6.0), Inches(3.5))
    new_box.fill.solid()
    new_box.fill.fore_color.rgb = RGBColor(30, 80, 50)
    new_box.line.color.rgb = GREEN

    new_title = slide.shapes.add_textbox(Inches(7.033), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = new_title.text_frame
    tf.paragraphs[0].text = "Agent 辅助方式"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GREEN

    new_content = slide.shapes.add_textbox(Inches(7.033), Inches(2.4), Inches(5.6), Inches(2.5))
    tf = new_content.text_frame
    tf.word_wrap = True
    items = ["/openresty-lua-plugins 5分钟", "/dev all (sync→build→test) 10分钟", "/fix-* (自动修复，如需要)", "", "总计: 15-30 分钟"]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

    # 效率提升
    efficiency_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.2))
    efficiency_box.fill.solid()
    efficiency_box.fill.fore_color.rgb = RGBColor(30, 80, 50)
    efficiency_box.line.color.rgb = GREEN

    eff_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.8))
    tf = eff_text.text_frame
    tf.paragraphs[0].text = "效率提升: 8-16 倍"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = YELLOW
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_roadmap_slide():
    """添加路线图幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "下一步计划"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 三个阶段
    stages = [
        ("当前状态", "Skills 工具集", "手动调用", CYAN),
        ("短期目标", "Agent 化", "自动调用\n智能路由\n上下文传递", YELLOW),
        ("长期愿景", "自主 Agent", "需求理解\n自主设计\n自主测试", GREEN)
    ]

    x_pos = 0.5
    for stage, title, desc, color in stages:
        # 阶段框
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.3), Inches(4.0), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(40, 40, 70)
        shape.line.color.rgb = color

        # 阶段名
        stage_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.5), Inches(3.6), Inches(0.4))
        tf = stage_box.text_frame
        tf.paragraphs[0].text = stage
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 标题
        title_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.9), Inches(3.6), Inches(0.5))
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 描述
        desc_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.5), Inches(3.6), Inches(1.2))
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        x_pos += 4.3

    # 箭头
    arrow1 = slide.shapes.add_textbox(Inches(4.3), Inches(2.3), Inches(0.5), Inches(0.5))
    tf = arrow1.text_frame
    tf.paragraphs[0].text = "→"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = YELLOW

    arrow2 = slide.shapes.add_textbox(Inches(8.6), Inches(2.3), Inches(0.5), Inches(0.5))
    tf = arrow2.text_frame
    tf.paragraphs[0].text = "→"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = YELLOW

    # 具体计划
    plan_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.2), Inches(12.333), Inches(2.5))
    plan_box.fill.solid()
    plan_box.fill.fore_color.rgb = RGBColor(30, 60, 50)
    plan_box.line.color.rgb = GREEN

    plan_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(11.9), Inches(0.5))
    tf = plan_title.text_frame
    tf.paragraphs[0].text = "具体计划"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GREEN

    plans = [
        "1. 技能编排自动化 - /dev all 自动串联修复链",
        "2. 知识库扩充 - 每次开发自动反馈到 Memory",
        "3. Agent 自主化 - 实现需求→交付闭环"
    ]

    plan_content = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.5))
    tf = plan_content.text_frame
    for i, plan in enumerate(plans):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = plan
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

    return slide


def add_token_optimization_slide():
    """添加Token优化幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Token 消耗优化策略"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 问题描述
    problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.5))
    tf = problem_box.text_frame
    p = tf.paragraphs[0]
    p.text = "问题：连续对话导致 Token 成本指数级上升"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY

    # 策略表
    strategies = [
        (".claudeignore", "忽略构建产物、依赖库、密钥", "减少 40% 无效加载"),
        ("CLAUDE.md 结构化", "明确模块介绍、目录结构", "避免重复理解"),
        ("/clear 清理", "开发不同模块时清理历史", "阻断指数增长"),
        ("/compact 压缩", "定期压缩对话历史", "保持上下文精简"),
        ("SKILL 化拆分", "细粒度技能 + Reference 按需加载", "精准 Token 消耗")
    ]

    y_pos = 1.6
    for strategy, method, effect in strategies:
        # 行背景
        row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.9))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = RGBColor(50, 50, 80)
        row_bg.line.color.rgb = RGBColor(80, 80, 100)

        # 策略名
        strat_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(3.0), Inches(0.7))
        tf = strat_box.text_frame
        tf.paragraphs[0].text = strategy
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = CYAN

        # 方法
        method_box = slide.shapes.add_textbox(Inches(3.8), Inches(y_pos + 0.1), Inches(5.0), Inches(0.7))
        tf = method_box.text_frame
        tf.paragraphs[0].text = method
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = WHITE

        # 效果
        effect_box = slide.shapes.add_textbox(Inches(9.0), Inches(y_pos + 0.1), Inches(3.6), Inches(0.7))
        tf = effect_box.text_frame
        tf.paragraphs[0].text = effect
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = YELLOW

        y_pos += 1.0

    return slide


def add_skill_design_slide():
    """添加Skill设计幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Skill 化 + 渐进式解析"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 核心思想
    idea_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.5))
    tf = idea_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心思想：细粒度拆分降低复杂度"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY

    # 对比
    # 传统方式
    old_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.0), Inches(2.0))
    old_box.fill.solid()
    old_box.fill.fore_color.rgb = RGBColor(80, 30, 30)
    old_box.line.color.rgb = RED

    old_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(5.6), Inches(0.4))
    tf = old_title.text_frame
    tf.paragraphs[0].text = "传统方式"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED

    old_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(5.6), Inches(1.2))
    tf = old_content.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "用户需求 → AI 理解全部项目"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "↓ 高 Token 消耗、高错误率"
    p.font.size = Pt(12)
    p.font.color.rgb = YELLOW
    p = tf.add_paragraph()
    p.text = "生成代码 → 可能出错"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Skill方式
    new_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.6), Inches(6.0), Inches(2.0))
    new_box.fill.solid()
    new_box.fill.fore_color.rgb = RGBColor(30, 80, 50)
    new_box.line.color.rgb = GREEN

    new_title = slide.shapes.add_textbox(Inches(7.033), Inches(1.8), Inches(5.6), Inches(0.4))
    tf = new_title.text_frame
    tf.paragraphs[0].text = "Skill 化方式"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GREEN

    new_content = slide.shapes.add_textbox(Inches(7.033), Inches(2.3), Inches(5.6), Inches(1.2))
    tf = new_content.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "用户需求 → 匹配 Skill → 加载 Reference"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "↓ 精准匹配、低 Token、高准确"
    p.font.size = Pt(12)
    p.font.color.rgb = YELLOW
    p = tf.add_paragraph()
    p.text = "生成代码 → 高准确率"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE

    # Reference 结构
    ref_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.9), Inches(12.333), Inches(3.0))
    ref_box.fill.solid()
    ref_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
    ref_box.line.color.rgb = CYAN

    ref_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.4))
    tf = ref_title.text_frame
    tf.paragraphs[0].text = "Reference 示例:"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CYAN

    ref_content = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(11.9), Inches(2.2))
    tf = ref_content.text_frame
    tf.word_wrap = True
    code = """.claude/skills/
├── openresty-lua-plugins/
│   ├── SKILL.md              # 技能定义（按需加载）
│   └── references/
│       ├── http.md           # HTTP 插件模板
│       └── stream.md         # TCP/UDP 插件模板
├── fix-compile/
│   └── SKILL.md              # 仅加载需要的修复规则
└── fix-test/
    └── SKILL.md"""
    tf.paragraphs[0].text = code
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = "Consolas"

    return slide


def add_e2e_slide():
    """添加E2E流程幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "E2E 流程打通"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN

    # 流水线
    pipeline_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Inches(1.8))
    pipeline_box.fill.solid()
    pipeline_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
    pipeline_box.line.color.rgb = CYAN

    # 流程步骤
    steps = ["sync", "build", "test", "analyze", "fix"]
    x_pos = 0.8
    for step in steps:
        step_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.5), Inches(2.0), Inches(0.8))
        step_shape.fill.solid()
        step_shape.fill.fore_color.rgb = CYAN
        step_shape.line.fill.background()

        tf = step_shape.text_frame
        tf.paragraphs[0].text = step
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK_BLUE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if step != "fix":
            arrow = slide.shapes.add_textbox(Inches(x_pos + 2.1), Inches(1.7), Inches(0.3), Inches(0.4))
            tf = arrow.text_frame
            tf.paragraphs[0].text = "→"
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.color.rgb = YELLOW

        x_pos += 2.4

    # 命令
    cmd_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(11.9), Inches(0.4))
    tf = cmd_box.text_frame
    tf.paragraphs[0].text = "/dev all 自动化流水线"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = YELLOW
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 失败处理
    fail_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.2), Inches(12.333), Inches(2.0))
    fail_box.fill.solid()
    fail_box.fill.fore_color.rgb = RGBColor(30, 60, 50)
    fail_box.line.color.rgb = GREEN

    fail_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(11.9), Inches(0.4))
    tf = fail_title.text_frame
    tf.paragraphs[0].text = "失败时自动处理流程"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GREEN

    fail_content = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(11.9), Inches(1.2))
    tf = fail_content.text_frame
    items = [
        "1. 匹配 tools/fixers/*.yaml 预制规则",
        "2. 输出结构化错误分析",
        "3. 建议调用 /fix-compile | /fix-test | /fix-runtime",
        "4. /fix-loop 自动迭代直到成功"
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

    # 跨平台
    cross_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.4), Inches(12.333), Inches(1.5))
    cross_box.fill.solid()
    cross_box.fill.fore_color.rgb = RGBColor(50, 50, 80)
    cross_box.line.color.rgb = YELLOW

    cross_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.4))
    tf = cross_title.text_frame
    tf.paragraphs[0].text = "跨平台透明化"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = YELLOW
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    cross_content = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.8))
    tf = cross_content.text_frame
    tf.paragraphs[0].text = "Windows 本地 → SSH (paramiko) → Linux 远程执行 → 返回结果"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_summary_slide():
    """添加总结幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "总结"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # 总结内容
    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(3.5))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(30, 50, 70)
    summary_box.line.color.rgb = YELLOW

    summary_text = """本项目基于 Google Agent 白皮书的 Model-Tools-Orchestration-Knowledge 架构理念，
构建了一套"项目规范驱动 + 技能细粒度拆分 + 知识自动沉淀 + 端到端流水线"的 AI 辅助开发体系：

• 以 CLAUDE.md 作为认知入口让 AI 理解项目上下文
• 以 Skills/ 按需加载 Reference 实现精准低耗调用
• 以 Memory/ 持续积累开发知识形成闭环
• 以 /dev all 串联 sync→build→test→fix 全流程自动化

最终将传统 4-6 小时的开发周期压缩至 15-30 分钟，
实现了从"人驱动工具"到"Agent 驱动开发"的范式跃迁。"""

    summary_content = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.9), Inches(3.2))
    tf = summary_content.text_frame
    tf.word_wrap = True
    for i, line in enumerate(summary_text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        if line.startswith('•'):
            p.font.color.rgb = CYAN

    # 效率提升
    eff_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.3), Inches(6.333), Inches(1.2))
    eff_box.fill.solid()
    eff_box.fill.fore_color.rgb = RGBColor(30, 80, 50)
    eff_box.line.color.rgb = GREEN

    eff_text = slide.shapes.add_textbox(Inches(3.5), Inches(5.5), Inches(6.333), Inches(0.8))
    tf = eff_text.text_frame
    tf.paragraphs[0].text = "效率提升: 8-16 倍"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = YELLOW
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_end_slide():
    """添加结束幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.alignment = PP_ALIGN.CENTER

    # Q&A
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.8))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Q & A"
    p.font.size = Pt(36)
    p.font.color.rgb = YELLOW
    p.alignment = PP_ALIGN.CENTER

    # 标签
    tags = ["OpenResty Skills", "Claude Code Agent", "端到端自动化开发"]
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.8))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "  |  ".join(tags)
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # 演示命令
    demo_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.5))
    tf = demo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "演示: /openresty-lua-plugins → /dev all → /feedback"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 120)
    p.alignment = PP_ALIGN.CENTER

    return slide


# 生成PPT
if __name__ == "__main__":
    # 封面
    add_title_slide("AI Coding 探索 Showcase", "基于 Claude Code 的 Agent 开发实践")

    # 内容页
    add_pain_point_slide()
    add_architecture_slide()
    add_solution_slide()
    add_comparison_slide()
    add_roadmap_slide()
    add_token_optimization_slide()
    add_skill_design_slide()
    add_e2e_slide()
    add_summary_slide()

    # 结束页
    add_end_slide()

    # 保存
    output_path = "D:/m30020610/WorkSpace/Skills/docs/ai-coding-showcase.pptx"
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    print(f"共 {len(prs.slides)} 页")